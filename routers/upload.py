# from fastapi import APIRouter, UploadFile, File, HTTPException
# import logging
# from models import UploadResponse
# from services.ingestion import process_document
# from services.embeddings import embed_texts
# from services.vector_store import upsert_chunks

# router = APIRouter(prefix="/upload", tags=["Document Upload"])
# logger = logging.getLogger(__name__)


# @router.post("/", response_model=UploadResponse, summary="Upload and index a PDF document")
# async def upload_document(file: UploadFile = File(...)):
#     """
#     Upload a PDF document to be processed and stored in the vector database.
#     """
#     try:
#         file_bytes = await file.read()
#         filename = file.filename or "unknown.pdf"

#         chunks, document_id = process_document(file_bytes, filename)

#         texts = [c["text"] for c in chunks]
#         embeddings = embed_texts(texts)

#         stored = upsert_chunks(chunks, embeddings)

#         logger.info(f"Upload complete: {filename} → {stored} chunks stored")
#         return UploadResponse(
#             message="Document processed and indexed successfully.",
#             document_id=document_id,
#             chunks_stored=stored,
#             filename=filename,
#         )

#     except ValueError as e:
#         raise HTTPException(status_code=422, detail=str(e))
#     except Exception as e:
#         logger.error(f"Upload failed: {e}", exc_info=True)
#         raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")



import os
import uuid
import logging
from fastapi import APIRouter, UploadFile, File, HTTPException
from models import UploadResponse
from services.embeddings import embed_texts
from services.vector_store import upsert_chunks

router = APIRouter(prefix="/upload", tags=["Document Upload"])
logger = logging.getLogger(__name__)

# ─────────────────────────────────────────
# EXTRACTORS
# ─────────────────────────────────────────

def extract_from_pdf(file_bytes: bytes) -> str:
    import fitz
    text = ""
    with fitz.open(stream=file_bytes, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text.strip()


def extract_from_docx(file_bytes: bytes) -> str:
    import io
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])


def extract_from_excel(file_bytes: bytes) -> str:
    import io
    import pandas as pd
    xls = pd.ExcelFile(io.BytesIO(file_bytes))
    all_text = []
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name)
        all_text.append(f"--- Sheet: {sheet_name} ---")
        all_text.append(df.to_string(index=False))
    return "\n\n".join(all_text)


def extract_from_csv(file_bytes: bytes) -> str:
    import io
    import pandas as pd
    df = pd.read_csv(io.BytesIO(file_bytes))
    return df.to_string(index=False)


def extract_from_pptx(file_bytes: bytes) -> str:
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    text = []
    for i, slide in enumerate(prs.slides, 1):
        text.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text.strip())
    return "\n".join(text)


def extract_from_text(file_bytes: bytes) -> str:
    for enc in ["utf-8", "latin-1", "cp1252"]:
        try:
            return file_bytes.decode(enc).strip()
        except UnicodeDecodeError:
            continue
    raise ValueError("Could not decode file with any known encoding.")


# ─────────────────────────────────────────
# TEXT EXTENSIONS
# ─────────────────────────────────────────

TEXT_EXTENSIONS = {
    ".txt", ".py", ".js", ".ts", ".jsx", ".tsx",
    ".java", ".cpp", ".c", ".cs", ".go", ".rs",
    ".html", ".css", ".json", ".xml", ".md",
    ".yaml", ".yml", ".env", ".sh", ".bat",
    ".sql", ".r", ".swift", ".kt", ".php", ".rb",
}


# ─────────────────────────────────────────
# MAIN EXTRACTOR
# ─────────────────────────────────────────

def extract_content(file_bytes: bytes, filename: str) -> tuple[str, str]:
    ext = os.path.splitext(filename)[1].lower()

    if ext == ".pdf":
        return extract_from_pdf(file_bytes), "pdf"
    elif ext in [".docx", ".doc"]:
        return extract_from_docx(file_bytes), "word"
    elif ext in [".xlsx", ".xls"]:
        return extract_from_excel(file_bytes), "excel"
    elif ext == ".csv":
        return extract_from_csv(file_bytes), "csv"
    elif ext in [".pptx", ".ppt"]:
        return extract_from_pptx(file_bytes), "powerpoint"
    elif ext in TEXT_EXTENSIONS:
        return extract_from_text(file_bytes), "code/text"
    else:
        # Last resort: try as plain text
        try:
            return extract_from_text(file_bytes), "unknown-text"
        except Exception:
            raise HTTPException(
                status_code=415,
                detail=f"Unsupported file type: '{ext}'"
            )


# ─────────────────────────────────────────
# CHUNKING
# ─────────────────────────────────────────

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    words = text.split()
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunks.append(" ".join(words[start:end]))
        start += chunk_size - overlap
    return chunks


# ─────────────────────────────────────────
# ENDPOINT
# ─────────────────────────────────────────

@router.post("/", response_model=UploadResponse, summary="Upload and index any file")
async def upload_document(file: UploadFile = File(...)):
    """
    Upload any file (PDF, Word, Excel, CSV, PowerPoint, code files, etc.)
    to be processed and stored in the vector database.
    """
    try:
        file_bytes = await file.read()
        filename = file.filename or "unknown"

        # Extract content based on file type
        extracted_text, file_type = extract_content(file_bytes, filename)

        if not extracted_text or len(extracted_text.strip()) < 5:
            raise HTTPException(
                status_code=400,
                detail="No content could be extracted from this file."
            )

        # Chunk the text
        document_id = str(uuid.uuid4())
        raw_chunks = chunk_text(extracted_text)

        # Build chunks in same format as your existing upsert_chunks expects
        chunks = [
            {
                "text": chunk,
                "document_id": document_id,
                "filename": filename,
                "chunk_index": i,
                "page": 0,
            }
            for i, chunk in enumerate(raw_chunks)
        ]

        # Embed
        texts = [c["text"] for c in chunks]
        embeddings = embed_texts(texts)

        # Store in Pinecone using your existing upsert_chunks
        stored = upsert_chunks(chunks, embeddings)

        logger.info(f"Upload complete: {filename} ({file_type}) → {stored} chunks stored")

        return UploadResponse(
            message=f"File processed and indexed successfully. Type: {file_type}",
            document_id=document_id,
            chunks_stored=stored,
            filename=filename,
        )

    except HTTPException:
        raise
    except ValueError as e:
        raise HTTPException(status_code=422, detail=str(e))
    except Exception as e:
        logger.error(f"Upload failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")